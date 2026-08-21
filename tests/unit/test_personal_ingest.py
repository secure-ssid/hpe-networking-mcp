"""Unit tests for cli_client.personal_ingest — local-only personal document
ingestion (PPTX/DOCX/PDF/VTT/text/HTML extraction, content-hash dedup,
chunk/embed/store, hybrid search).

All fixtures are synthetic and built in ``tmp_path``. No test references any
real user document path — this module must never leak the existence or
contents of a real local folder into the committed test suite.

Covers two real defects found while ingesting a genuine internal-docs
collection (fixed in the same change, regression-tested here):

- ``extract_pptx_text`` crashed with ``AttributeError`` when a slide reports
  ``has_notes_slide=True`` but its ``notes_text_frame`` is ``None`` (a
  malformed/edited notes placeholder) — one bad slide used to blank out an
  entire multi-hundred-slide deck's extracted text.
- ``extract_pptx_text`` (via ``Presentation()``) raised ``BadZipFile`` for a
  single corrupt embedded media stream (bad CRC-32), blocking extraction of
  an otherwise-intact deck. ``_rebuild_pptx_without_corrupt_members`` repairs
  this by replacing only the broken zip member(s) with an empty placeholder.
- ``_chunk_text`` (mirroring ``ingestion/chunking.py``) used to leave a short
  heading immediately followed by a blank line and an oversized body
  paragraph as its own tiny, low-information "orphan" chunk instead of
  folding it into a neighboring chunk.
"""

from __future__ import annotations

import io
import zipfile
import zlib
from pathlib import Path

import pytest

from hpe_networking_mcp.cli_client import personal_ingest as pi

# ── chunking ──────────────────────────────────────────────────────────────


def test_chunk_text_short_text_single_chunk():
    chunks = pi._chunk_text("A short sentence.")
    assert chunks == ["A short sentence."]


def test_chunk_text_long_text_multiple_bounded_chunks():
    long_text = "This is a repeated sentence about CX switches. " * 60
    chunks = pi._chunk_text(long_text)
    assert len(chunks) > 1
    assert all(len(c) <= pi._CHUNK_SIZE for c in chunks)


def test_chunk_text_merges_orphan_heading_into_following_body():
    # Mirrors ingestion/chunking.py's regression test for the same
    # real-world defect (a short heading immediately followed by a blank
    # line and an oversized body paragraph must not survive as its own
    # tiny, low-information chunk).
    heading = "# Widget X specifications"
    body_paragraph_1 = "x" * 790
    body_paragraph_2 = "y" * 250
    text = f"{heading}\n\n{body_paragraph_1}\n\n{body_paragraph_2}"

    chunks = pi._chunk_text(text)

    assert heading not in chunks
    assert all(len(chunk) >= pi._MIN_CHUNK_SIZE for chunk in chunks)
    assert any(heading in chunk and "x" in chunk for chunk in chunks)
    assert any("y" * 250 in chunk for chunk in chunks)


def test_merge_small_chunks_folds_trailing_chunk_backward():
    chunks = ["a" * 300, "b" * 300, "tiny tail"]

    merged = pi._merge_small_chunks(chunks)

    assert merged == ["a" * 300, "b" * 300 + "\n\ntiny tail"]


# ── PPTX extraction ───────────────────────────────────────────────────────


def _build_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "CX 6300 Overview"
    body = slide.placeholders[1]
    body.text_frame.text = "16p 1G 8p Smart Rate 10G Class8 PoE"

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker notes about the 6300 family"

    # second slide with a table
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    rows, cols = 2, 2
    table_shape = slide2.shapes.add_table(rows, cols, 0, 0, 2000000, 500000)
    table = table_shape.table
    table.cell(0, 0).text = "SKU"
    table.cell(0, 1).text = "Ports"
    table.cell(1, 0).text = "S6X58A"
    table.cell(1, 1).text = "24"

    prs.save(str(path))


def test_extract_pptx_text_slides_notes_and_tables(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    _build_pptx(path)

    text = pi.extract_pptx_text(path)

    assert "## Slide 1" in text
    assert "CX 6300 Overview" in text
    assert "Class8 PoE" in text
    assert "Notes: Speaker notes about the 6300 family" in text
    assert "## Slide 2" in text
    assert "SKU | Ports" in text
    assert "S6X58A | 24" in text


def test_extract_pptx_text_survives_none_notes_text_frame(tmp_path: Path):
    """Regression: a notes slide with has_notes_slide=True but a missing body
    placeholder (notes_text_frame is None) must not crash extraction, and the
    slide's other text must still be captured."""
    path = tmp_path / "deck_bad_notes.pptx"
    _build_pptx(path)

    from pptx import Presentation

    prs = Presentation(str(path))
    slide = prs.slides[0]
    notes_slide = slide.notes_slide
    assert notes_slide.notes_text_frame is not None  # sanity before corruption
    body_ph = notes_slide.placeholders[1]
    body_ph._element.getparent().remove(body_ph._element)
    assert notes_slide.notes_text_frame is None  # confirms the corruption landed
    prs.save(str(path))

    text = pi.extract_pptx_text(path)  # must not raise

    assert "CX 6300 Overview" in text
    assert "Class8 PoE" in text
    assert "Notes:" not in text  # no notes captured, but nothing else lost
    assert "## Slide 2" in text  # later slides still processed


def test_extract_pptx_text_recovers_from_corrupt_zip_member(tmp_path: Path):
    """Regression: a single embedded-media zip member with a bad CRC-32 must
    not block extraction of an otherwise-valid deck — exactly the real-world
    failure (a corrupt embedded ``.m4a`` audio blob) this fix targets. Only
    non-structural (non-``.xml``/``.rels``) members are safe to zero out, so
    the corrupted member here is a synthetic ``ppt/media/`` blob, appended
    with ``ZIP_STORED`` so a single flipped byte deterministically produces
    a clean CRC-32 mismatch rather than an unpredictable DEFLATE failure."""
    path = tmp_path / "deck_corrupt_media.pptx"
    _build_pptx(path)

    media_name = "ppt/media/media99.bin"
    with zipfile.ZipFile(path, "a") as zf:
        zf.writestr(
            zipfile.ZipInfo(media_name),
            b"fake embedded audio bytes" * 20,
            zipfile.ZIP_STORED,
        )

    raw = bytearray(path.read_bytes())
    with zipfile.ZipFile(io.BytesIO(bytes(raw))) as zf:
        info = zf.getinfo(media_name)
        offset = info.header_offset + 30 + len(info.filename) + len(info.extra)
    raw[offset] ^= 0xFF  # flip a byte inside the (uncompressed) media data
    path.write_bytes(bytes(raw))

    with pytest.raises(zipfile.BadZipFile):
        zipfile.ZipFile(path).read(media_name)

    text = pi.extract_pptx_text(path)  # must not raise, must recover

    assert "## Slide 1" in text
    assert "CX 6300 Overview" in text
    assert "## Slide 2" in text


def test_extract_pptx_text_raises_when_structural_part_is_corrupt(tmp_path: Path):
    """A corrupt *structural* OPC part (``.xml``/``.rels``) is genuinely not
    recoverable by zeroing — that would only trade one unparseable package
    for another. This must surface as a clean, catchable exception (handled
    by ingest_folder's per-file try/except) rather than a confusing deep
    lxml/XML parser error."""
    path = tmp_path / "deck_corrupt_structural.pptx"
    _build_pptx(path)

    raw = bytearray(path.read_bytes())
    with zipfile.ZipFile(io.BytesIO(bytes(raw))) as zf:
        target = next(
            i for i in zf.infolist() if i.filename.endswith(".xml") and i.compress_size > 0
        )
        offset = target.header_offset + 30 + len(target.filename) + len(target.extra)
    raw[offset] ^= 0xFF
    path.write_bytes(bytes(raw))

    with pytest.raises(Exception):  # noqa: B017 - deliberately broad: any load failure is acceptable
        pi.extract_pptx_text(path)


def test_rebuild_pptx_without_corrupt_members_keeps_good_entries(tmp_path: Path):
    path = tmp_path / "mini.zip"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("good.txt", b"kept content")
        zf.writestr("bad.bin", b"will be corrupted")

    raw = bytearray(path.read_bytes())
    with zipfile.ZipFile(io.BytesIO(bytes(raw))) as zf:
        bad_info = zf.getinfo("bad.bin")
        offset = bad_info.header_offset + 30 + len(bad_info.filename) + len(bad_info.extra)
    raw[offset] ^= 0xFF
    path.write_bytes(bytes(raw))

    rebuilt = pi._rebuild_pptx_without_corrupt_members(path)
    with zipfile.ZipFile(rebuilt) as zf:
        assert zf.read("good.txt") == b"kept content"
        assert zf.read("bad.bin") == b""  # corrupt member replaced, not left broken


def test_rebuild_pptx_without_corrupt_members_handles_zlib_error(tmp_path: Path):
    """A broken DEFLATE bitstream raises a raw ``zlib.error`` (NOT a
    ``zipfile.BadZipFile`` subclass — confirmed via MRO inspection) instead
    of a clean CRC mismatch. The rebuild must catch and recover from this
    too, not just ``BadZipFile``."""
    path = tmp_path / "mini_deflate.zip"
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("good.txt", b"kept content")
        # highly compressible payload so the DEFLATE stream has internal
        # back-references a corrupted byte can break ("invalid distance too
        # far back") rather than merely flipping an output byte.
        zf.writestr("bad.bin", b"abcdefgh" * 200)

    raw = bytearray(path.read_bytes())
    with zipfile.ZipFile(io.BytesIO(bytes(raw))) as zf:
        bad_info = zf.getinfo("bad.bin")
        offset = bad_info.header_offset + 30 + len(bad_info.filename) + len(bad_info.extra)
    # Offset +10 into this exact deterministic DEFLATE stream is confirmed
    # (empirically, via direct probing) to reliably break the decompressor's
    # internal back-reference state rather than merely flip an output byte.
    raw[offset + 10] ^= 0xFF
    path.write_bytes(bytes(raw))

    with zipfile.ZipFile(path) as zf:
        try:
            zf.read("bad.bin")
            pytest.skip("byte flip happened not to corrupt this zlib stream")
        except zipfile.BadZipFile:
            pytest.skip("byte flip produced a clean CRC mismatch, not zlib.error")
        except zlib.error:
            pass  # exactly the case under test

    rebuilt = pi._rebuild_pptx_without_corrupt_members(path)
    with zipfile.ZipFile(rebuilt) as zf:
        assert zf.read("good.txt") == b"kept content"
        assert zf.read("bad.bin") == b""


def test_rebuild_pptx_without_corrupt_members_reraises_for_structural_part(tmp_path: Path):
    """A corrupt ``.xml``/``.rels`` member is NOT safe to zero out — that
    would silently trade one broken package for another. The rebuild must
    refuse (raise) rather than return an equally-unusable result."""
    path = tmp_path / "mini_structural.zip"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("good.txt", b"kept content")
        zf.writestr("[Content_Types].xml", b"<Types>corrupt me</Types>")

    raw = bytearray(path.read_bytes())
    with zipfile.ZipFile(io.BytesIO(bytes(raw))) as zf:
        bad_info = zf.getinfo("[Content_Types].xml")
        offset = bad_info.header_offset + 30 + len(bad_info.filename) + len(bad_info.extra)
    raw[offset] ^= 0xFF
    path.write_bytes(bytes(raw))

    with pytest.raises(zipfile.BadZipFile):
        pi._rebuild_pptx_without_corrupt_members(path)


# ── DOCX extraction ───────────────────────────────────────────────────────


def test_extract_docx_text_paragraphs_and_tables(tmp_path: Path):
    import docx

    path = tmp_path / "doc.docx"
    document = docx.Document()
    document.add_paragraph("CX 6300M switch family overview.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "SKU"
    table.cell(0, 1).text = "Ports"
    table.cell(1, 0).text = "S6X59A"
    table.cell(1, 1).text = "48"
    document.save(str(path))

    text = pi.extract_docx_text(path)

    assert "CX 6300M switch family overview." in text
    assert "SKU | Ports" in text
    assert "S6X59A | 48" in text


# ── PDF extraction ────────────────────────────────────────────────────────


def _build_pdf(path: Path, text: str) -> None:
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    page = writer.add_blank_page(width=200, height=200)

    stream_obj = DecodedStreamObject()
    stream_obj.set_data(f"BT /F1 14 Tf 10 100 Td ({text}) Tj ET".encode())
    content_ref = writer._add_object(stream_obj)

    font_dict = DictionaryObject()
    font_dict[NameObject("/Type")] = NameObject("/Font")
    font_dict[NameObject("/Subtype")] = NameObject("/Type1")
    font_dict[NameObject("/BaseFont")] = NameObject("/Helvetica")
    font_ref = writer._add_object(font_dict)

    font_res = DictionaryObject()
    font_res[NameObject("/F1")] = font_ref
    resources = DictionaryObject()
    resources[NameObject("/Font")] = font_res

    page[NameObject("/Resources")] = resources
    page[NameObject("/Contents")] = content_ref

    with path.open("wb") as fh:
        writer.write(fh)


def test_extract_pdf_text_returns_page_text(tmp_path: Path):
    path = tmp_path / "sheet.pdf"
    _build_pdf(path, "CX6300 datasheet contents")

    text = pi.extract_pdf_text(path)

    assert "CX6300 datasheet contents" in text


# ── VTT extraction ────────────────────────────────────────────────────────


def test_extract_vtt_text_dedupes_rolling_captions(tmp_path: Path):
    path = tmp_path / "captions.vtt"
    path.write_text(
        "WEBVTT\n\n"
        "1\n00:00:00.000 --> 00:00:02.000\nHello and\n\n"
        "2\n00:00:02.000 --> 00:00:04.000\nHello and welcome\n\n"
        "3\n00:00:04.000 --> 00:00:06.000\nHello and welcome to the demo\n\n",
        encoding="utf-8",
    )

    text = pi.extract_vtt_text(path)

    # Rolling captions collapse to only the final, fullest form.
    assert text == "Hello and welcome to the demo"


def test_extract_vtt_text_keeps_distinct_cues(tmp_path: Path):
    path = tmp_path / "captions2.vtt"
    path.write_text(
        "WEBVTT\n\n"
        "1\n00:00:00.000 --> 00:00:02.000\nFirst distinct sentence.\n\n"
        "2\n00:00:02.000 --> 00:00:04.000\nSecond unrelated sentence.\n\n",
        encoding="utf-8",
    )

    text = pi.extract_vtt_text(path)

    assert "First distinct sentence." in text
    assert "Second unrelated sentence." in text
    assert text.count("\n") == 1  # two distinct lines, no rolling collapse


# ── suffix dispatch ───────────────────────────────────────────────────────


def test_extract_text_dispatches_markdown(tmp_path: Path):
    path = tmp_path / "note.md"
    path.write_text("# Title\n\nBody text.", encoding="utf-8")
    assert pi.extract_text(path) == "# Title\n\nBody text."


def test_extract_text_returns_none_for_unsupported_suffix(tmp_path: Path):
    path = tmp_path / "sheet.xlsm"
    path.write_bytes(b"not a real xlsm")
    assert pi.extract_text(path) is None


def test_extract_text_returns_none_on_extraction_error(tmp_path: Path):
    path = tmp_path / "broken.pdf"
    path.write_bytes(b"not a real pdf")
    assert pi.extract_text(path) is None  # extract_pdf_text raises; caught -> None


# ── collection name validation ───────────────────────────────────────────


def test_validate_collection_accepts_lowercase_alnum_underscore():
    assert pi._validate_collection("internal_docs_1") == "internal_docs_1"


@pytest.mark.parametrize("bad", ["Internal", "internal-docs", "internal docs", "internal;drop"])
def test_validate_collection_rejects_invalid_names(bad: str):
    with pytest.raises(ValueError):
        pi._validate_collection(bad)


# ── full ingest_folder pipeline ──────────────────────────────────────────


def test_ingest_folder_extracts_chunks_and_is_searchable(tmp_path: Path):
    docs_dir = tmp_path / "docs"
    docs_dir.mkdir()
    (docs_dir / "cx6300.md").write_text(
        "# CX 6300M Switch\n\nThe CX 6300M supports 90W Class8 PoE and MACsec.",
        encoding="utf-8",
    )
    (docs_dir / "mx301.txt").write_text(
        "MX301 is a Juniper routing platform for self-driving networks.",
        encoding="utf-8",
    )
    data_dir = tmp_path / "personal-lancedb"

    result = pi.ingest_folder(docs_dir, collection="internal", data_dir=data_dir)

    assert result.files_seen == 2
    assert result.files_ingested == 2
    assert result.files_failed == 0
    assert result.chunks_written >= 2

    counts = pi.personal_collection_counts(data_dir=data_dir)
    assert counts.get("internal") == result.chunks_written

    hits = pi.search_personal("CX 6300M PoE", collection="internal", data_dir=data_dir)
    assert hits
    assert any("6300" in h.get("text", "") for h in hits)


def test_ingest_folder_is_idempotent_on_rerun(tmp_path: Path):
    docs_dir = tmp_path / "docs"
    docs_dir.mkdir()
    (docs_dir / "note.md").write_text("Idempotent ingestion note.", encoding="utf-8")
    data_dir = tmp_path / "personal-lancedb"

    first = pi.ingest_folder(docs_dir, collection="internal", data_dir=data_dir)
    assert first.files_ingested == 1

    second = pi.ingest_folder(docs_dir, collection="internal", data_dir=data_dir)
    assert second.files_ingested == 0
    assert second.files_skipped_unchanged == 1


def test_ingest_folder_dedupes_identical_content_under_different_names(tmp_path: Path):
    docs_dir = tmp_path / "docs"
    docs_dir.mkdir()
    content = "Duplicate content saved under two different file names."
    (docs_dir / "original.md").write_text(content, encoding="utf-8")
    (docs_dir / "original (1).md").write_text(content, encoding="utf-8")
    data_dir = tmp_path / "personal-lancedb"

    result = pi.ingest_folder(docs_dir, collection="internal", data_dir=data_dir)

    assert result.files_seen == 2
    assert result.files_ingested == 1
    assert result.files_skipped_duplicate == 1


def test_ingest_folder_reports_unsupported_and_continues(tmp_path: Path):
    docs_dir = tmp_path / "docs"
    docs_dir.mkdir()
    (docs_dir / "keep.md").write_text("Readable content.", encoding="utf-8")
    (docs_dir / "calc.xlsm").write_bytes(b"binary spreadsheet content")
    data_dir = tmp_path / "personal-lancedb"

    result = pi.ingest_folder(docs_dir, collection="internal", data_dir=data_dir)

    assert result.files_seen == 2
    assert result.files_ingested == 1
    assert result.files_skipped_unsupported == 1
    assert result.files_failed == 0


def test_search_personal_returns_empty_list_when_no_index_exists(tmp_path: Path):
    data_dir = tmp_path / "never-created-lancedb"
    assert pi.search_personal("anything", collection="internal", data_dir=data_dir) == []


def test_personal_collection_counts_empty_when_no_index_exists(tmp_path: Path):
    data_dir = tmp_path / "never-created-lancedb"
    assert pi.personal_collection_counts(data_dir=data_dir) == {}


def test_default_personal_data_dir_is_outside_repo():
    repo_root = Path(__file__).resolve().parents[2]
    resolved = pi.default_personal_data_dir()
    assert repo_root not in resolved.parents
    assert ".config" in resolved.parts
    assert "hpe-mcp" in resolved.parts
    assert "personal" in resolved.parts
