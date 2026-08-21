"""Local, content-searchable personal document ingestion.

Extends ``documents.py``'s metadata-only ``DocumentStore`` with real text
extraction (PPTX, DOCX, PDF, VTT, Markdown/HTML/text) plus a self-contained
chunk/embed/LanceDB pipeline (``langchain_text_splitters`` +
``pipeline.clients.embed_client`` + ``pipeline.clients.lance_client``).

Everything here is stored under a **personal, non-repository** LanceDB
directory (default ``~/.config/hpe-mcp/personal/lancedb``) — never the
repository's committed/gitignored ``data/`` directory, and never uploaded
anywhere. Embedding runs fully in-process via fastembed (ONNX); no document
text or extracted content leaves the local machine. This is the right place
for personal/internal/proprietary material (e.g. internal sales/technical
enablement decks) that must stay local and must never be bundled into the
shared, redistributable corpus under ``data/docs.lance``.
"""

from __future__ import annotations

import hashlib
import io
import re
import uuid
import zipfile
import zlib
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable

# NOTE: chunking is intentionally reimplemented here (matching
# ``ingestion/chunking.py``'s CHUNK_SIZE/CHUNK_OVERLAP/separators exactly)
# rather than imported from the top-level ``ingestion`` package. That package
# lives outside ``src/`` and is only importable when the process has the
# repo root on ``sys.path`` (e.g. ``uv run python -c`` or pytest's rootdir
# insertion). The real MCP router is launched as
# ``python3 src/hpe_networking_mcp/mcp_servers/tool_router.py`` with
# ``PYTHONPATH=<repo>/src`` only (see .cursor/mcp.json) — under that exact
# launch, ``sys.path[0]`` is the script's own directory, not the repo root,
# so ``from ingestion.chunking import chunk_text`` raises ModuleNotFoundError
# (reproduced directly). ``langchain-text-splitters`` is already a base
# dependency importable from anywhere, so this module depends on it directly
# instead of on the fragile top-level package.
from langchain_text_splitters import RecursiveCharacterTextSplitter

from hpe_networking_mcp.cli_client.config import default_user_data_dir

_CHUNK_SIZE = 800
_CHUNK_OVERLAP = 100
#: Chunks shorter than this are folded into a neighboring chunk instead of
#: being embedded/indexed standalone (mirrors ``ingestion/chunking.py``'s
#: MIN_CHUNK_SIZE -- see ``_merge_small_chunks`` there for the full
#: rationale: a short heading immediately followed by a blank line and an
#: oversized body paragraph is otherwise left as its own tiny, low-content
#: "orphan" chunk that can outrank the real content for title-keyword
#: queries).
_MIN_CHUNK_SIZE = 200
_splitter = RecursiveCharacterTextSplitter(
    chunk_size=_CHUNK_SIZE,
    chunk_overlap=_CHUNK_OVERLAP,
    separators=["\n\n", "\n", ". ", " ", ""],
)


def _merge_small_chunks(chunks: list[str]) -> list[str]:
    """Fold chunks under ``_MIN_CHUNK_SIZE`` into an adjacent chunk.

    Kept in sync with ``ingestion/chunking.py``'s ``_merge_small_chunks``
    (duplicated rather than imported for the same ``sys.path`` reason
    documented above for chunking itself).
    """
    if len(chunks) <= 1:
        return chunks

    merged: list[str] = []
    pending = ""
    for chunk in chunks:
        candidate = f"{pending}\n\n{chunk}" if pending else chunk
        if len(candidate) < _MIN_CHUNK_SIZE:
            pending = candidate
            continue
        merged.append(candidate)
        pending = ""

    if pending:
        if merged:
            merged[-1] = f"{merged[-1]}\n\n{pending}"
        else:
            merged.append(pending)

    return merged


def _chunk_text(text: str) -> list[str]:
    return _merge_small_chunks(_splitter.split_text(text))


SUPPORTED_SUFFIXES = {
    ".md",
    ".markdown",
    ".txt",
    ".htm",
    ".html",
    ".pdf",
    ".pptx",
    ".docx",
    ".vtt",
}


def default_personal_data_dir() -> Path:
    """Personal LanceDB root — outside the repo, outside git, never shared."""
    return default_user_data_dir() / "personal" / "lancedb"


def _md5_uuid(key: str) -> str:
    """Deterministic row id from a stable key -- not a security/cryptographic use."""
    return str(uuid.UUID(hashlib.md5(key.encode(), usedforsecurity=False).hexdigest()))


def _content_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for block in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


# ── Text extraction ──────────────────────────────────────────────────────────


def _rebuild_pptx_without_corrupt_members(path: Path) -> io.BytesIO:
    """Rebuild ``path`` in memory with corrupt, non-structural zip member(s)
    replaced by an empty placeholder. Only called as a fallback after
    python-pptx has already failed to open the file directly — the
    expensive full-archive scan below is never paid on the common,
    well-formed-file path.

    A corrupt member can fail in two different ways: a clean
    ``zipfile.BadZipFile`` (decompression succeeds but the CRC-32 doesn't
    match) or a raw ``zlib.error`` (the DEFLATE bitstream itself is broken,
    e.g. "invalid distance too far back") — the latter is NOT a
    ``BadZipFile`` subclass, so both must be caught here.

    Only members that are safe to discard (i.e. not ``.xml``/``.rels`` — the
    OPC package structure, content-type map, and slide XML) are zeroed out.
    Real-world corruption is essentially always in an embedded binary media
    blob (audio/video/image); zeroing a corrupt *structural* part would only
    trade one unrecoverable error for another (an unparseable empty XML
    document), so if one of those is corrupt we re-raise the original error
    instead of returning a package that still can't load.
    """
    with zipfile.ZipFile(path) as zin:
        infos = zin.infolist()
        bad_names = set()
        critical_bad_names = set()
        for info in infos:
            try:
                zin.read(info.filename)
            except (zipfile.BadZipFile, zlib.error, OSError):
                lower = info.filename.lower()
                if lower.endswith((".xml", ".rels")):
                    critical_bad_names.add(info.filename)
                else:
                    bad_names.add(info.filename)
        if critical_bad_names:
            raise zipfile.BadZipFile(
                "corrupt structural OPC part(s), not recoverable: "
                + ", ".join(sorted(critical_bad_names))
            )
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for info in infos:
                if info.filename in bad_names:
                    zout.writestr(info.filename, b"")
                else:
                    zout.writestr(info, zin.read(info.filename))
    buf.seek(0)
    return buf


def extract_pptx_text(path: Path) -> str:
    """Slide shape text + speaker notes, in slide order.

    Defensive by design: a single malformed shape/slide (e.g. a notes
    placeholder with no ``notes_text_frame``, seen in real internal decks)
    must not sacrifice the rest of a multi-hundred-slide deck's content, and
    neither must a corrupt embedded media stream. python-opc (python-pptx's
    package loader) eagerly reads every zip member — including embedded
    video/audio blobs — when opening a .pptx, so a single corrupt embedded
    media stream (bad CRC-32, or a broken DEFLATE stream raising a raw
    ``zlib.error``) otherwise blocks the *entire* file even though the slide
    XML/text content is completely intact. On the first such failure we
    rebuild an in-memory copy with just the broken member(s) zeroed out and
    retry once.
    """
    from pptx import Presentation

    try:
        prs = Presentation(str(path))
    except (zipfile.BadZipFile, zlib.error):
        prs = Presentation(_rebuild_pptx_without_corrupt_members(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        try:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    lines.append(shape.text_frame.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            lines.append(" | ".join(cells))
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                notes = notes_frame.text.strip() if notes_frame is not None else ""
                if notes:
                    lines.append(f"Notes: {notes}")
        except AttributeError:
            # Malformed shape/notes XML on this one slide — keep whatever
            # text this slide already collected and move on to the rest.
            pass
        if lines:
            parts.append(f"## Slide {i}\n" + "\n".join(lines))
    return "\n\n".join(parts)



def extract_docx_text(path: Path) -> str:
    """Paragraph text + table cell text, in document order."""
    import docx

    doc = docx.Document(str(path))
    parts: list[str] = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def extract_pdf_text(path: Path) -> str:
    """Per-page text, same pypdf pattern used by ingestion/scrape_security_lifecycle.py."""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p.strip() for p in pages if p.strip())


_VTT_TIMESTAMP_RE = re.compile(
    r"^\d{2}:\d{2}:\d{2}[.,]\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}[.,]\d{3}"
)
_VTT_CUE_NUM_RE = re.compile(r"^\d+$")
_VTT_TAG_RE = re.compile(r"<[^>]+>")


def extract_vtt_text(path: Path) -> str:
    """Caption text only — strip the WEBVTT header, cue numbers, timestamps,
    and markup tags, then collapse rolling-caption duplicates.

    Auto-generated captions commonly emit each cue as a growing rolling
    sentence (cue 2's text is a superset of cue 1's), which would otherwise
    multiply the transcript several-fold. Keep a cue's text only when it is
    not already a substring of the next kept cue.
    """
    raw = path.read_text(encoding="utf-8", errors="ignore")
    cues: list[str] = []
    buffer: list[str] = []
    for line in raw.splitlines():
        stripped = line.strip()
        if not stripped:
            if buffer:
                cues.append(" ".join(buffer))
                buffer = []
            continue
        if stripped.upper().startswith("WEBVTT"):
            continue
        if stripped.upper().startswith(("NOTE", "STYLE", "REGION")):
            continue
        if _VTT_TIMESTAMP_RE.match(stripped):
            continue
        if _VTT_CUE_NUM_RE.match(stripped):
            continue
        buffer.append(_VTT_TAG_RE.sub("", stripped))
    if buffer:
        cues.append(" ".join(buffer))

    deduped: list[str] = []
    for cue in cues:
        cue = cue.strip()
        if not cue:
            continue
        if deduped and cue.startswith(deduped[-1]):
            deduped[-1] = cue  # rolling caption grew — replace with the fuller version
            continue
        if deduped and deduped[-1].startswith(cue):
            continue  # this cue is a strict prefix of what we already kept
        deduped.append(cue)
    return "\n".join(deduped)


def _extract_text_dispatch(path: Path) -> str | None:
    """Dispatch by suffix. Raises on extraction failure (caller decides how
    to report it); returns None only for a suffix with no extractor."""
    suffix = path.suffix.lower()
    if suffix in (".md", ".markdown", ".txt"):
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix in (".htm", ".html"):
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
        return soup.get_text("\n")
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    if suffix == ".docx":
        return extract_docx_text(path)
    if suffix == ".vtt":
        return extract_vtt_text(path)
    return None


def extract_text(path: Path) -> str | None:
    """Dispatch by suffix. Returns None for unsupported or unreadable files."""
    try:
        return _extract_text_dispatch(path)
    except Exception:  # noqa: BLE001 - best-effort single-file preview helper
        return None


# ── Ingestion (chunk + embed + store) ────────────────────────────────────────


@dataclass
class IngestResult:
    files_seen: int = 0
    files_ingested: int = 0
    files_skipped_unchanged: int = 0
    files_skipped_duplicate: int = 0
    files_skipped_unsupported: int = 0
    files_failed: int = 0
    chunks_written: int = 0
    errors: dict[str, str] | None = None


_COLLECTION_RE = re.compile(r"^[a-z0-9_]+$")


def _validate_collection(collection: str) -> str:
    if not _COLLECTION_RE.match(collection):
        raise ValueError(
            f"invalid collection name: {collection!r} (lowercase letters/digits/underscore only)"
        )
    return collection


def _stamped_hash(stamped_key: str) -> str:
    return stamped_key.rsplit("::", 1)[-1]


def _existing_file_hashes(db, table_name: str, collection: str) -> set[str]:
    """Content hashes already ingested for this collection, regardless of
    which file path they were ingested under (dedupes identical files that
    exist under multiple names/copies)."""
    from hpe_networking_mcp.pipeline.clients import lance_client

    table = lance_client.docs_table(db, table_name)
    if table is None:
        return set()
    arrow = (
        table.search()
        .select(["file_path"])
        .where(f"source = '{_validate_collection(collection)}'", prefilter=True)
        .limit(table.count_rows())
        .to_arrow()
    )
    return {_stamped_hash(p) for p in arrow.column("file_path").to_pylist()}


def ingest_folder(
    folder: Path,
    *,
    collection: str = "internal",
    data_dir: Path | None = None,
    recursive: bool = True,
    progress: Callable[[str], None] | None = None,
) -> IngestResult:
    """Extract, chunk, embed, and store every supported file under ``folder``.

    Idempotent by content: a file whose SHA-256 already exists in this
    collection (under any file name/path — catches identical files saved as
    multiple copies) is skipped without re-extracting or re-embedding. Safe
    to re-run to pick up only new/changed files.
    """
    from hpe_networking_mcp.pipeline.clients import lance_client
    from hpe_networking_mcp.pipeline.clients.embed_client import EmbedClient

    _validate_collection(collection)
    folder = Path(folder).expanduser().resolve()
    data_dir = Path(data_dir) if data_dir else default_personal_data_dir()
    result = IngestResult()
    errors: dict[str, str] = {}

    pattern = "**/*" if recursive else "*"
    files = sorted(
        p for p in folder.glob(pattern) if p.is_file() and not p.name.startswith(".")
    )

    db = lance_client.connect(data_dir)
    existing_hashes = _existing_file_hashes(db, lance_client.DOCS_TABLE, collection)
    processed_this_run: set[str] = set()

    all_records: list[dict[str, Any]] = []
    all_texts: list[str] = []

    for path in files:
        result.files_seen += 1
        if path.suffix.lower() not in SUPPORTED_SUFFIXES:
            result.files_skipped_unsupported += 1
            continue

        file_path_key = str(path)
        file_hash = _file_hash(path)
        stamped_key = f"{file_path_key}::{file_hash}"
        if file_hash in existing_hashes:
            result.files_skipped_unchanged += 1
            continue
        if file_hash in processed_this_run:
            result.files_skipped_duplicate += 1
            if progress:
                progress(f"SKIP (duplicate content already ingested this run): {path.name}")
            continue

        try:
            text = _extract_text_dispatch(path)
        except Exception as exc:  # noqa: BLE001 - one bad file must not abort the batch
            result.files_failed += 1
            errors[file_path_key] = f"{type(exc).__name__}: {exc}"
            if progress:
                progress(f"FAILED ({type(exc).__name__}): {path.name}")
            continue
        if not text or not text.strip():
            result.files_failed += 1
            errors[file_path_key] = "no extractable text"
            if progress:
                progress(f"SKIP (no extractable text): {path.name}")
            continue

        chunks = _chunk_text(text)
        for i, chunk in enumerate(chunks):
            record_id = _md5_uuid(f"{collection}:{file_hash}:{i}")
            all_records.append(
                {
                    "id": record_id,
                    "text": chunk,
                    "source": collection,
                    "doc_type": "personal_internal",
                    "file_path": stamped_key,
                    "chunk_index": i,
                    "content_hash": _content_hash(chunk),
                }
            )
            all_texts.append(chunk)
        processed_this_run.add(file_hash)
        result.files_ingested += 1
        result.chunks_written += len(chunks)
        if progress:
            progress(f"extracted {len(chunks)} chunks: {path.name}")

    if all_records:
        embedder = EmbedClient()
        vectors = embedder.embed_document(all_texts)
        for record, vector in zip(all_records, vectors):
            record["vector"] = vector
        lance_client.merge_docs_rows(db, all_records, lance_client.DOCS_TABLE)
        table = lance_client.docs_table(db, lance_client.DOCS_TABLE)
        if table is not None:
            lance_client.build_fts_index(table)

    result.errors = errors
    return result


def search_personal(
    query: str,
    *,
    collection: str | None = None,
    top_k: int = 10,
    data_dir: Path | None = None,
) -> list[dict[str, Any]]:
    """Hybrid search over the personal LanceDB store (same shape as rag.search_docs)."""
    from hpe_networking_mcp.pipeline.clients import lance_client
    from hpe_networking_mcp.pipeline.clients.embed_client import EmbedClient

    data_dir = Path(data_dir) if data_dir else default_personal_data_dir()
    db = lance_client.connect(data_dir)
    embedder = EmbedClient()
    vector = embedder.embed_query(query)
    try:
        return lance_client.hybrid_search(
            db, query, vector, top_k=top_k, source_filter=collection
        )
    except FileNotFoundError:
        return []


def personal_collection_counts(data_dir: Path | None = None) -> dict[str, int]:
    from hpe_networking_mcp.pipeline.clients import lance_client

    data_dir = Path(data_dir) if data_dir else default_personal_data_dir()
    db = lance_client.connect(data_dir)
    return lance_client.source_counts(db, lance_client.DOCS_TABLE)
